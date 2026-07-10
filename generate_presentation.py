from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

DARK_BG = RGBColor(0x1E, 0x2A, 0x3D)
ACCENT = RGBColor(0x00, 0xA8, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
TABLE_HEADER_BG = RGBColor(0x00, 0x86, 0xBB)
TABLE_ROW_ALT = RGBColor(0x2A, 0x3A, 0x50)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    if width is None:
        width = prs.slide_width
    if height is None:
        height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=18, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_accent_bar(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def make_slide(title_text, subtitle_text=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.9), title_text, font_size=36, bold=True, color=ACCENT)
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5), subtitle_text, font_size=18, color=LIGHT_GRAY)
    return slide

# ---------- SLIDE 1: Title ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), prs.slide_height)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
    "SQLTraine", font_size=56, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
    "Интерактивный SQL-тренажёр", font_size=32, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.6),
    "Повышайте навыки написания SQL-запросов с автоматической проверкой", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
    "Django 6 · React 18 · PostgreSQL 15 · Docker", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ---------- SLIDE 2: About ----------
slide = make_slide("О проекте")

items = [
    "▎Web-приложение для отработки навыков написания SQL-запросов",
    "▎Автоматическая проверка решений — сравнение с эталонным результатом",
    "▎Целевая аудитория: студенты, начинающие разработчики, преподаватели",
    "",
    "▎Возможности:",
    "  • Регистрация и авторизация (JWT)",
    "  • Библиотека задач по категориям и сложности",
    "  • Встроенный редактор кода с подсветкой синтаксиса",
    "  • Статьи с теорией и навигацией",
    "  • История попыток и статистика прогресса",
    "  • Тёмная и светлая темы",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), items, font_size=20)

# ---------- SLIDE 3: Tech Stack ----------
slide = make_slide("Технологический стек")

stack_data = [
    ("Компонент", "Технология"),
    ("Backend", "Django 6 + Django REST Framework"),
    ("Frontend", "React 18 + TypeScript"),
    ("База данных", "PostgreSQL 15+"),
    ("Редактор кода", "Monaco Editor"),
    ("Стилизация", "Tailwind CSS"),
    ("Аутентификация", "JWT (djangorestframework-simplejwt)"),
    ("Контейнеризация", "Docker + Docker Compose"),
]

rows, cols = len(stack_data), 2
table = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.8), Inches(10), Inches(4.8)).table
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(6.5)

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = stack_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.name = "Calibri"
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BG

# ---------- SLIDE 4: Architecture ----------
slide = make_slide("Архитектура приложения")

items = [
    "▎Три сервиса в Docker Compose, объединённые сетью sqltrain:",
    "",
    "  ┌─────────────┐      ┌──────────────┐      ┌─────────────┐",
    "  │  Frontend   │ ──→  │   Backend    │ ──→  │  Database   │",
    "  │ React + Vite│      │  Django DRF  │      │ PostgreSQL  │",
    "  │  port 5173  │      │  port 8000   │      │  port 5432  │",
    "  └─────────────┘      └──────────────┘      └─────────────┘",
    "",
    "▎Frontend: React (Vite) с Monaco Editor для написания SQL",
    "▎Backend: Django 6 + DRF — REST API, проверка запросов, бизнес-логика",
    "▎Database: PostgreSQL — хранение задач, пользователей, попыток",
    "▎Песочница: отдельное подключение к БД для изолированного выполнения SQL",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), items, font_size=18)

# ---------- SLIDE 5: Roles ----------
slide = make_slide("Ролевая модель и аутентификация")

items = [
    "▎Три роли пользователей:",
    "",
    "  • Student (Студент) — решение задач, просмотр теории, статистика",
    "  • Teacher (Преподаватель) — создание и управление задачами",
    "  • Admin (Администратор) — полный доступ, управление пользователями",
    "",
    "▎Аутентификация: JWT-токены",
    "  • Access Token — 2 часа",
    "  • Refresh Token — 7 дней",
    "  • Bearer-схема авторизации",
    "",
    "▎Эндпоинты:",
    "  • POST /api/auth/register/   — регистрация",
    "  • POST /api/auth/login/      — вход",
    "  • POST /api/auth/refresh/    — обновление токена",
    "  • GET  /api/auth/me/         — текущий пользователь",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), items, font_size=18)

# ---------- SLIDE 6: Features 1 ----------
slide = make_slide("Функциональность — Управление задачами")

items = [
    "▎CRUD-операции с задачами (для Teacher/Admin)",
    "",
    "▎Поля задачи:",
    "  • title — название задачи",
    "  • description — условие (Markdown)",
    "  • difficulty — сложность (Easy, Medium, Hard)",
    "  • category — категория (SELECT, JOIN, GROUP BY, Subqueries, DDL, DML)",
    "  • initial_sql — исходная схема БД и тестовые данные",
    "  • expected_result — эталонный результат (JSON)",
    "  • hints — подсказки",
    "  • test_cases — дополнительные тест-кейсы",
    "  • verification_query — SELECT для проверки DDL/DML задач",
    "",
    "▎Фильтрация задач по категории и сложности",
    "▎Публикация/скрытие задач",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), items, font_size=18)

# ---------- SLIDE 7: Features 2 ----------
slide = make_slide("Функциональность — Решение задач")

items = [
    "▎Встроенный редактор Monaco с подсветкой SQL-синтаксиса",
    "",
    "▎Процесс решения:",
    "  1. Пользователь выбирает задачу из списка",
    "  2. Изучает условие и исходную схему БД",
    "  3. Пишет SQL-запрос в редакторе",
    "  4. Нажимает «Выполнить» — видит результат запроса",
    "  5. Нажимает «Отправить» — система проверяет решение",
    "  6. Получает обратную связь: правильно / ошибка",
    "",
    "▎Автоматическая проверка:",
    "  • Сравнение результата с expected_result",
    "  • Проверка по дополнительным test_cases",
    "  • Для DDL/DML — verification_query после запроса пользователя",
    "",
    "▎Подсказки: заготовленные заранее или через ИИ",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), items, font_size=18)

# ---------- SLIDE 8: Features 3 ----------
slide = make_slide("Функциональность — Теория и профиль")

items = [
    "▎Статьи с теорией по SQL:",
    "  • Список статей с поиском",
    "  • Навигационный бар внутри статьи",
    "  • Статьи связаны с задачами (related_articles)",
    "",
    "▎Профиль пользователя:",
    "  • Имя и фото профиля",
    "  • Количество решённых задач",
    "  • Статистика и прогресс",
    "  • Дней без пропуска (streak)",
    "",
    "▎Дополнительно:",
    "  • Комментарии под задачами",
    "  • Тёмная / светлая тема",
    "  • Таблица лидеров",
    "  • Сценарные подборки задач (tutorial)",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), items, font_size=18)

# ---------- SLIDE 9: API ----------
slide = make_slide("API Эндпоинты")

api_data = [
    ("Метод", "URL", "Описание"),
    ("POST", "/api/auth/register/", "Регистрация"),
    ("POST", "/api/auth/login/", "Вход"),
    ("POST", "/api/auth/refresh/", "Обновление токена"),
    ("GET", "/api/auth/me/", "Текущий пользователь"),
    ("GET", "/api/tasks/", "Список задач"),
    ("POST", "/api/tasks/", "Создание задачи"),
    ("GET", "/api/tasks/{id}/", "Детали задачи"),
    ("PUT", "/api/tasks/{id}/", "Обновление задачи"),
    ("DELETE", "/api/tasks/{id}/", "Удаление задачи"),
    ("POST", "/api/tasks/{id}/execute/", "Выполнить запрос"),
    ("POST", "/api/tasks/{id}/submit/", "Отправить на проверку"),
    ("GET", "/api/progress/", "Прогресс пользователя"),
    ("GET", "/api/submissions/", "История попыток"),
]

rows, cols = len(api_data), 3
table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2)).table
table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(4.2)
table.columns[2].width = Inches(5.7)

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = api_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.name = "Calibri"
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BG

# ---------- SLIDE 10: DB ----------
slide = make_slide("Структура базы данных")

items = [
    "▎Модели Django ORM:",
    "",
    "  • User (встроенная модель Django + роли)",
    "    — id, username, email, password_hash, role (student/teacher/admin)",
    "",
    "  • Category",
    "    — id, name, description",
    "",
    "  • Task",
    "    — id, name, description, difficulty, category, schema_sql",
    "    — expected_query, hints, verification_query, is_published",
    "    — related_articles (ManyToMany)",
    "",
    "  • Submission",
    "    — id, user, task, user_query, is_correct, error_message",
    "",
    "  • Article (приложение theory)",
    "    — статьи с теорией, связанные с задачами",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), items, font_size=17)

# ---------- SLIDE 11: Deployment ----------
slide = make_slide("Развёртывание")

items = [
    "▎Docker Compose — 3 сервиса:",
    "",
    "  ┌─────────────────────────────────────────────────┐",
    "  │  docker-compose.yml                             │",
    "  │                                                 │",
    "  │  db:       postgres:latest   (port 5433:5432)   │",
    "  │  backend:  Django + Gunicorn (port 8000)        │",
    "  │  frontend: React + Vite      (port 5173)        │",
    "  └─────────────────────────────────────────────────┘",
    "",
    "▎Healthcheck для PostgreSQL — pg_isready",
    "▎Backend зависит от db (condition: service_healthy)",
    "▎Frontend зависит от backend",
    "▎Hot-reload для frontend через volume (./src:/app/src)",
    "",
    "▎Запуск одной командой:",
    "  docker compose up --build",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), items, font_size=18)

# ---------- SLIDE 12: Thank You ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), prs.slide_height)

add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
    "Спасибо за внимание!", font_size=48, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
    "SQLTraine — Интерактивный SQL-тренажёр", font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
    "Репозиторий: github.com/your-username/SQLTraine", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ---------- Save ----------
output_path = "SQLTraine_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved as {output_path}")
